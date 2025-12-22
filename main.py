from importlib.resources import contents
from fastapi import FastAPI, UploadFile, File, Form, types
from fastapi.middleware.cors import CORSMiddleware
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import csv
import pathlib
import textwrap

try:
    from IPython.display import display, Markdown
except Exception:
    # IPython is optional for server runs; provide no-op fallbacks so
    # importing this module doesn't fail when IPython isn't installed.
    def display(*_args, **_kwargs):
        return None

    def Markdown(text):
        return text
import io
from pydantic import BaseModel

import vertexai
from vertexai.generative_models import GenerativeModel
from vertexai.generative_models import (
    GenerativeModel,
    Part,
)
# from google.generativeai.types import Part


# Initialize Vertex AI ONCE (startup time)
vertexai.init(
    project="poc-script-genai",
    location="us-central1",
)

model = GenerativeModel(
    "projects/poc-script-genai/locations/us-central1/endpoints/1201184567508074496"
)


app = FastAPI()

# CORS React frontend support
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:5173",
        "http://127.0.0.1:5173",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# Schemas
class ChatRequest(BaseModel):
    message: str


@app.get("/")
def root():
    return {"status": "bkl chalrha hai"}

# @app.post("/chat")
# def chat(req: ChatRequest):
#     response=model.generate_content(req.message)  
#     return {"reply": response.text}



SYSTEM_INSTRUCTIONS = """ 
Untill not asked to create/generate a script, you will act as a normal chatbot.
You are a content writer, when asked to create/generate a script, you will follow the structure below exactly, UNTILL then you are instructed otherwise.0
Based on the provided inputs, generate a complete script and all associated metadata. The output must follow the specified structure and format.

## Output Format
A table in MarkDown format with three columns. 
Each table row MUST be on a new line.
Do NOT collapse multiple rows into a single line.
First column is Time in seconds, second column is Voice Over text, third column is Visuals description.
| Time (s) | Voice Over                                    | Visuals                                                                 |
| :------- | :---------------------------------------------- | :---------------------------------------------------------------------- |

## Output Structure

### Title
A compelling title for the video.

### Description
A detailed description for the video platform (e.g., YouTube).

### Details
- Video Length: Total length of the video in seconds.
- Word Count: Total number of words in the Voice Over.

### Script
A 3-column Markdown table with the following headers: `Time (s)`, `Voice Over`, `Visuals`.



## Constraints
- The `Time (s)` column must be in cumulative seconds.
- The `Voice Over` and `Visuals` columns should contain concise, clear sentences.

## Example

### Title
How to Brew the Perfect Pour-Over Coffee in Under 3 Minutes

### Description
Learn the art of brewing the perfect pour-over coffee in just under 3 minutes! This step-by-step guide will walk you through the process, from selecting the right beans to mastering your pouring technique

### Details
- Video Length: 180 seconds
- Word Count: 150 words

### Script
| Time (s) | Voice Over                                    | Visuals                                                                 |
| :------- | :---------------------------------------------- | :---------------------------------------------------------------------- |

"""

@app.post("/chat")
async def chat(
    prompt: str = Form(""),
    file: UploadFile | None = File(None)
):
    parts = []
    parts.append(SYSTEM_INSTRUCTIONS)

    # ---------- PDF ----------
    if file and file.filename.lower().endswith(".pdf"):
        pdf_bytes = await file.read()
        reader = PdfReader(io.BytesIO(pdf_bytes))

        pdf_text = ""
        for page in reader.pages:
            if page.extract_text():
                pdf_text += page.extract_text() + "\n"

        parts.append(f"""
DOCUMENT (PDF):
{pdf_text}

USER TASK:
{prompt}
""")

    # ---------- IMAGE ----------
    elif file and file.filename.lower().endswith((".png", ".jpg", ".jpeg")):
        image_bytes = await file.read()
        parts.append(Part.from_data(image_bytes, mime_type=file.content_type))
        parts.append(prompt)

    # ---------- VIDEO ----------
    elif file and file.filename.lower().endswith(".mp4"):
        video_bytes = await file.read()
        parts.append(
            Part.from_data(video_bytes, mime_type="video/mp4")
        )
        parts.append(prompt)

    # ---------- TEXT ----------
    elif file and file.filename.lower().endswith(".txt"):
        text_bytes = await file.read()
        text_content = text_bytes.decode("utf-8", errors="ignore")

        parts.append(f"""
DOCUMENT (TEXT):
{text_content}

USER TASK:
{prompt}
""")

    # ---------- CSV ----------
    elif file and file.filename.lower().endswith(".csv"):
        csv_bytes = await file.read()
        decoded = csv_bytes.decode("utf-8", errors="ignore")

        reader = csv.reader(io.StringIO(decoded))
        csv_text = "\n".join([", ".join(row) for row in reader])

        parts.append(f"""
DOCUMENT (CSV):
{csv_text}

USER TASK:
{prompt}
""")

    # ---------- DOCX ----------
    elif file and file.filename.lower().endswith(".docx"):
        doc_bytes = await file.read()
        doc = Document(io.BytesIO(doc_bytes))

        doc_text = "\n".join([para.text for para in doc.paragraphs])

        parts.append(f"""
DOCUMENT (DOCX):
{doc_text}

USER TASK:
{prompt}
""")

    # ---------- XLSX ----------
    elif file and file.filename.lower().endswith(".xlsx"):
        xlsx_bytes = await file.read()
        wb = load_workbook(io.BytesIO(xlsx_bytes), data_only=True)

        excel_text = ""
        for sheet in wb:
            excel_text += f"\nSHEET: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                excel_text += ", ".join([str(cell) if cell else "" for cell in row]) + "\n"

        parts.append(f"""
DOCUMENT (EXCEL):
{excel_text}

USER TASK:
{prompt}
""")

    # ---------- PPTX ----------
    elif file and file.filename.lower().endswith(".pptx"):
        ppt_bytes = await file.read()
        prs = Presentation(io.BytesIO(ppt_bytes))

        ppt_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    ppt_text += shape.text + "\n"

        parts.append(f"""
DOCUMENT (PRESENTATION):
{ppt_text}

USER TASK:
{prompt}
""")

    # ---------- TEXT ONLY ----------
    else:
        parts.append(prompt)


    # response = model.generate_content(parts)
    # return {"reply": response.text}
    final_text = ""

    stream = model.generate_content(
        parts,
        stream=True
    )

    for response in stream:
        if not response.candidates:
            continue

        candidate = response.candidates[0]

        if not candidate.content or not candidate.content.parts:
            continue

        for part in candidate.content.parts:
            if hasattr(part, "text") and part.text:
                final_text += part.text

    def sanitize_output(text: str) -> str:
        replacements = {
            "###": "",
            "***": "",
            "**": "",
        }
        for k, v in replacements.items():
            text = text.replace(k, v)
        return text.strip()

    def normalize_markdown_table(text: str) -> str:
        """
        Converts pipe-collapsed markdown into a proper row-based markdown table.
        """

        # Remove accidental double pipes and trim
        text = text.replace("||", "|").replace("**", "").strip()

        # Split by pipe and clean cells
        cells = [c.strip() for c in text.split("||") if c.strip()]

        # Expected columns: 3
        cols = 3

        # If not divisible, return original text
        if len(cells) < cols:
            return text

        rows = [cells[i:i+cols] for i in range(0, len(cells), cols)]

        # Build markdown table
        table = []
        table.append(f"| {rows[0][0]} | {rows[0][1]} | {rows[0][2]} |")
        table.append("| :------- | :--------- | :--------- |")

        for row in rows[1:]:
            if len(row) == cols:
                table.append(f"| {row[0]} | {row[1]} | {row[2]} |")

        return "\n".join(table)

    # ✅ POST-PROCESSING PIPELINE
    clean_text = sanitize_output(final_text)
    markdown_table = normalize_markdown_table(clean_text)
    return {"reply": markdown_table}





# i am giving a speech on LLMs , i have attached a document go through it and create a 5 minute speech. target audience are students who are AI enthusiast . keep the tone of the speech more interesting and interactive, try to add examples , idioms and jokes whereever posible