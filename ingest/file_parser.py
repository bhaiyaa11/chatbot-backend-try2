# import io
# import csv
# import asyncio
# import logging
# from fastapi import UploadFile
# from pypdf import PdfReader
# from docx import Document
# from openpyxl import load_workbook
# from pptx import Presentation
# from pdf2image import convert_from_bytes
# import pytesseract
# from vertexai.generative_models import Part
# from config import TOKEN_BUDGETS

# logger = logging.getLogger(__name__)

# from pypdf import PdfReader
# import io

# def extract_pdf_text(pdf_bytes: bytes) -> str:
#     reader = PdfReader(io.BytesIO(pdf_bytes))
#     text = []

#     for page in reader.pages:
#         content = page.extract_text()
#         if content:
#             text.append(content)

#     return "\n".join(text)
# # ───────────────────────────────────────────────────────────────
# # Helpers
# # ───────────────────────────────────────────────────────────────

# def _trim_files(texts: list[str], total_budget: int) -> list[str]:
#     """Distribute token budget evenly across text files."""
#     if not texts:
#         return []

#     per_file = max(total_budget // len(texts), 1000)
#     return [t[:per_file] for t in texts]


# def _is_meaningful(text: str) -> bool:
#     """
#     Detect whether extracted PDF text is real content
#     or garbage from image-based PDFs.
#     """
#     if not text:
#         return False

#     words = text.split()
#     # heuristic threshold
#     return len(words) > 150


# # ───────────────────────────────────────────────────────────────
# # Main API
# # ───────────────────────────────────────────────────────────────

# async def parse_files(files: list[UploadFile], stage: str = "VOICE_OVER") -> list:
#     """
#     Parse uploaded files concurrently.
#     Returns mixture of:
#         - text strings
#         - Vertex AI Part objects (images/videos)
#     """

#     if not files:
#         return []

#     results = await asyncio.gather(*[_parse_one(f) for f in files])

#     # ---- flatten results (important for PDF image lists) ----
#     flat_results = []
#     for r in results:
#         if isinstance(r, list):
#             flat_results.extend(r)
#         else:
#             flat_results.append(r)

#     # Separate text vs media
#     text_parts = [r for r in flat_results if isinstance(r, str) and r.strip()]
#     media_parts = [r for r in flat_results if not isinstance(r, str) and r is not None]

#     # Apply token budget only to text
#     budget = TOKEN_BUDGETS.get(stage, {}).get("file_budget", 15_000)
#     trimmed_texts = _trim_files(text_parts, budget)

#     return trimmed_texts + media_parts


# # ───────────────────────────────────────────────────────────────
# # Router
# # ───────────────────────────────────────────────────────────────

# async def _parse_one(file: UploadFile):
#     name = file.filename.lower()
#     data = await file.read()

#     try:
#         if name.endswith(".pdf"):
#             return await asyncio.to_thread(_pdf, data)

#         elif name.endswith(".docx"):
#             return await asyncio.to_thread(_docx, data)

#         elif name.endswith(".xlsx"):
#             return await asyncio.to_thread(_xlsx, data)

#         elif name.endswith(".pptx"):
#             return await asyncio.to_thread(_pptx, data)

#         elif name.endswith(".csv"):
#             return await asyncio.to_thread(_csv, data)

#         elif name.endswith(".txt"):
#             return data.decode(errors="ignore")

#         elif name.endswith((".png", ".jpg", ".jpeg", ".webp")):
#             return Part.from_data(
#                 data=data,
#                 mime_type=file.content_type or "image/png"
#             )

#         elif name.endswith((".mp4", ".mov", ".webm", ".mkv")):
#             return Part.from_data(
#                 data=data,
#                 mime_type=file.content_type or "video/mp4"
#             )

#         else:
#             logger.info(f"Unsupported file skipped: {file.filename}")
#             return None

#     except Exception as e:
#         logger.error(f"Parse error [{name}]: {e}")
#         return None


# # ───────────────────────────────────────────────────────────────
# # PDF PARSER (FIXED)
# # ───────────────────────────────────────────────────────────────

# def _pdf(data: bytes):
#     """
#     Smart PDF parser:

#     1. Try text extraction
#     2. Validate text quality
#     3. If poor → treat as image PDF
#     4. Send pages as images to multimodal model
#     """

#     reader = PdfReader(io.BytesIO(data))

#     extracted = "\n".join(
#         page.extract_text() or ""
#         for page in reader.pages[:15]
#     )

#     # ✅ GOOD TEXT PDF
#     if _is_meaningful(extracted):
#         logger.info("PDF parsed as TEXT document")
#         return extracted

#     # ✅ IMAGE PDF (slides / diagrams / scans)
#     logger.info("PDF detected as IMAGE document → converting pages to images")

#     images = convert_from_bytes(
#         data,
#         dpi=350,          # higher quality for diagrams
#         last_page=min(5, len(reader.pages))
#     )

#     parts = []

#     for img in images:
#         buffer = io.BytesIO()
#         img.save(buffer, format="PNG")

#         parts.append(
#             Part.from_data(
#                 data=buffer.getvalue(),
#                 mime_type="image/png"
#             )
#         )

#     return parts


# # ───────────────────────────────────────────────────────────────
# # Other File Parsers
# # ───────────────────────────────────────────────────────────────

# def _docx(data: bytes) -> str:
#     doc = Document(io.BytesIO(data))
#     return "\n".join(p.text for p in doc.paragraphs)


# def _xlsx(data: bytes) -> str:
#     wb = load_workbook(io.BytesIO(data), data_only=True)

#     lines = []
#     for sheet in wb:
#         for row in sheet.iter_rows(values_only=True):
#             lines.append(", ".join(str(c or "") for c in row))

#     return "\n".join(lines)


# def _pptx(data: bytes) -> str:
#     prs = Presentation(io.BytesIO(data))

#     texts = []
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 texts.append(shape.text)

#     return "\n".join(texts)


# def _csv(data: bytes) -> str:
#     decoded = data.decode(errors="ignore")
#     reader = csv.reader(io.StringIO(decoded))
#     return "\n".join(", ".join(r) for r in reader)






import io
import csv
import asyncio
import logging
from typing import Union

from fastapi import UploadFile
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from google.genai import types

from config import TOKEN_BUDGETS

logger = logging.getLogger(__name__)

__all__ = ["parse_files", "extract_pdf_text"]


# ───────────────────────────────────────────────────────────────
# PDF TEXT EXTRACTION
# ───────────────────────────────────────────────────────────────

def parse_pdf(file_bytes: bytes) -> str:
    """Production-safe PDF text extraction."""
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        text: list[str] = []

        for page in reader.pages:
            content = page.extract_text()
            if content:
                text.append(content)

        return "\n".join(text)

    except Exception as e:  # noqa: BLE001
        logger.error("PDF text extraction failed: %s", e, exc_info=True)
        return ""


def extract_pdf_text(pdf_bytes: bytes) -> str:
    """Public helper wrapper (kept for compatibility)."""
    return parse_pdf(pdf_bytes)


# ───────────────────────────────────────────────────────────────
# Private helpers
# ───────────────────────────────────────────────────────────────

def _trim_files(texts: list[str], total_budget: int) -> list[str]:
    """Distribute token budget evenly across text files."""
    if not texts:
        return []
    per_file = max(total_budget // len(texts), 1000)
    return [t[:per_file] for t in texts]


def _is_meaningful(text: str) -> bool:
    """
    Detect whether extracted PDF text is real content
    or garbage from an image-based PDF.
    """
    if not text:
        return False
    return len(text.split()) > 150


# ───────────────────────────────────────────────────────────────
# Main API
# ───────────────────────────────────────────────────────────────

async def parse_files(files: list[UploadFile], stage: str = "VOICE_OVER") -> list:
    """
    Parse uploaded files concurrently.

    Returns:
        - text strings
        - Vertex AI Part objects (images/videos)
    """
    if not files:
        return []

    results = await asyncio.gather(*[_parse_one(f) for f in files])

    # Flatten results
    flat_results: list = []
    for r in results:
        if isinstance(r, list):
            flat_results.extend(r)
        else:
            flat_results.append(r)

    # Separate text vs media
    text_parts: list[str] = [
        r for r in flat_results if isinstance(r, str) and r.strip()
    ]

    media_parts: list[types.Part] = [
        r for r in flat_results if not isinstance(r, str) and r is not None
    ]

    # Apply token budget only to text
    budget: int = TOKEN_BUDGETS.get(stage, {}).get("file_budget", 15_000)
    trimmed_texts = _trim_files(text_parts, budget)

    return trimmed_texts + media_parts


# ───────────────────────────────────────────────────────────────
# Router
# ───────────────────────────────────────────────────────────────

async def _parse_one(file: UploadFile) -> Union[str, types.Part, list, None]:
    name = (file.filename or "").lower()
    data = await file.read()

    try:
        if name.endswith(".pdf"):
            return await asyncio.to_thread(_pdf, data)

        if name.endswith(".docx"):
            return await asyncio.to_thread(_docx, data)

        if name.endswith(".xlsx"):
            return await asyncio.to_thread(_xlsx, data)

        if name.endswith(".pptx"):
            return await asyncio.to_thread(_pptx, data)

        if name.endswith(".csv"):
            return await asyncio.to_thread(_csv, data)

        if name.endswith(".txt"):
            return data.decode(errors="ignore")

        if name.endswith((".png", ".jpg", ".jpeg", ".webp")):
            return types.Part.from_bytes(
                data=data,
                mime_type=file.content_type or "image/png",
            )

        if name.endswith((".mp4", ".mov", ".webm", ".mkv")):
            return types.Part.from_bytes(
                data=data,
                mime_type=file.content_type or "video/mp4",
            )

        logger.info("Unsupported file skipped: %s", file.filename)
        return None

    except Exception as e:  # noqa: BLE001
        logger.error("Parse error [%s]: %s", name, e, exc_info=True)
        return None


# ───────────────────────────────────────────────────────────────
# PDF parser (FUNCTIONALITY UNCHANGED)
# ───────────────────────────────────────────────────────────────

def _pdf(data: bytes) -> Union[str, list[types.Part]]:
    """
    Smart PDF parser:

    1. Try text extraction.
    2. Validate text quality.
    3. If poor → treat as image PDF.
    4. Return pages as Vertex AI image Parts.
    """

    extracted = parse_pdf(data)

    if _is_meaningful(extracted):
        logger.info("PDF parsed as TEXT document")
        return extracted

    logger.info("PDF detected as IMAGE document → attempting image conversion")

    # ✅ PRODUCTION FIX: lazy import (prevents Vercel crash)
    try:
        from pdf2image import convert_from_bytes
    except ImportError:
        logger.warning(
            "pdf2image not installed — skipping image conversion fallback"
        )
        return extracted

    reader = PdfReader(io.BytesIO(data))

    images = convert_from_bytes(
        data,
        dpi=350,
        first_page=1,
        last_page=min(5, len(reader.pages)),
    )

    parts: list[types.Part] = []

    for img in images:
        buffer = io.BytesIO()
        img.save(buffer, format="PNG")

        parts.append(
            types.Part.from_bytes(
                data=buffer.getvalue(),
                mime_type="image/png",
            )
        )

    return parts


# ───────────────────────────────────────────────────────────────
# Other file parsers
# ───────────────────────────────────────────────────────────────

def _docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs)


def _xlsx(data: bytes) -> str:
    wb = load_workbook(io.BytesIO(data), data_only=True)
    lines: list[str] = []

    for sheet in wb:
        for row in sheet.iter_rows(values_only=True):
            lines.append(", ".join(str(c or "") for c in row))

    return "\n".join(lines)


def _pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    texts: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)

    return "\n".join(texts)


def _csv(data: bytes) -> str:
    decoded = data.decode(errors="ignore")
    reader = csv.reader(io.StringIO(decoded))
    return "\n".join(", ".join(r) for r in reader)